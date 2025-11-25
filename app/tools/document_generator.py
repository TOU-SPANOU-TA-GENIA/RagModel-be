# app/tools/document_generator.py
"""
Document generation tool - creates files based on user specifications.
Supports Word (.docx), PowerPoint (.pptx), PDF, and other formats.
"""

from typing import Dict, Any, Optional, List
from pathlib import Path
from dataclasses import dataclass
import json

from app.tools.base import BaseTool, ToolResult
from app.logger import setup_logger
from app.config import BASE_DIR

logger = setup_logger(__name__)


@dataclass
class DocumentSpec:
    """Specification for document generation."""
    doc_type: str  # docx, pptx, pdf, txt, md
    title: str
    content: Any  # Structure depends on doc_type
    metadata: Dict[str, Any]


class DocumentGeneratorTool(BaseTool):
    """
    Generates documents from structured specifications.
    Delegates to specific generators based on file type.
    """
    
    def __init__(self, output_dir: Path = None):
        super().__init__()
        self.output_dir = output_dir or (BASE_DIR / "outputs")
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Registry of generators
        self.generators = {
            'docx': self._generate_word,
            'pptx': self._generate_powerpoint,
            'pdf': self._generate_pdf,
            'txt': self._generate_text,
            'md': self._generate_markdown,
        }
    
    @property
    def name(self) -> str:
        return "generate_document"
    
    @property
    def description(self) -> str:
        return """Generate documents (Word, PowerPoint, PDF, etc.) based on specifications.
        Accepts: doc_type, title, content structure.
        Returns: file path for download."""
    
    # app/tools/document_generator.py - Update _execute_impl()

    def _execute_impl(self, doc_type: str, title: str, 
                    content: Any, **kwargs) -> ToolResult:
        """Generate document based on type."""
        
        doc_type = doc_type.lower().strip('.')
        
        if doc_type not in self.generators:
            return ToolResult(
                success=False,
                error=f"Unsupported document type: {doc_type}",
                data=None
            )
        
        try:
            # If content is just a string (the query), structure it
            if isinstance(content, str) and len(content) > 50:
                content = self._structure_content(content, doc_type, title)
            
            # Generate the document
            generator = self.generators[doc_type]
            file_path = generator(title, content, **kwargs)
            
            return ToolResult(
                success=True,
                data={
                    "file_path": str(file_path),
                    "file_name": file_path.name,
                    "file_type": doc_type,
                    "file_size_bytes": file_path.stat().st_size,
                    "title": title,
                    "download_ready": True
                }
            )
            
        except Exception as e:
            logger.error(f"Document generation failed: {e}", exc_info=True)
            return ToolResult(
                success=False,
                error=f"Failed to generate {doc_type}: {str(e)}",
                data=None
            )

    def _structure_content(self, query: str, doc_type: str, title: str) -> Dict:
        """Extract structured content from query."""
        import re
        
        # Simple extraction - look for sections mentioned
        sections = {}
        
        # Pattern: "covering X, Y, and Z"
        match = re.search(r'covering\s+(.+)', query, re.IGNORECASE)
        if match:
            topics = match.group(1)
            # Split by commas and 'and'
            items = re.split(r',\s*(?:and\s+)?|\s+and\s+', topics)
            
            for item in items:
                item = item.strip()
                if item:
                    section_name = item.capitalize()
                    # Generate some placeholder content
                    sections[section_name] = [
                        f"Point 1 about {item}",
                        f"Point 2 about {item}",
                        f"Point 3 about {item}"
                    ]
        
        # If no sections found, create generic structure
        if not sections:
            sections = {
                "Overview": [title],
                "Main Content": ["Content related to " + title],
                "Conclusion": ["Summary"]
            }
        
        return sections
    
    
    def _generate_word(self, title: str, content: Any, **kwargs) -> Path:
        """Generate Word document."""
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
        except ImportError:
            raise ImportError("python-docx required. Install: pip install python-docx")
        
        doc = Document()
        
        # Add title
        heading = doc.add_heading(title, level=0)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Process content structure
        if isinstance(content, dict):
            self._add_structured_content_to_word(doc, content)
        elif isinstance(content, list):
            for item in content:
                if isinstance(item, dict):
                    self._add_section_to_word(doc, item)
                else:
                    doc.add_paragraph(str(item))
        else:
            doc.add_paragraph(str(content))
        
        # Save
        file_path = self.output_dir / f"{self._sanitize_filename(title)}.docx"
        doc.save(str(file_path))
        
        logger.info(f"Generated Word document: {file_path}")
        return file_path
    
    def _add_structured_content_to_word(self, doc, content: Dict):
        """Add structured content to Word document."""
        for section_title, section_content in content.items():
            doc.add_heading(section_title, level=1)
            
            if isinstance(section_content, list):
                for item in section_content:
                    doc.add_paragraph(str(item), style='List Bullet')
            elif isinstance(section_content, dict):
                self._add_structured_content_to_word(doc, section_content)
            else:
                doc.add_paragraph(str(section_content))
    
    def _add_section_to_word(self, doc, section: Dict):
        """Add a section to Word document."""
        if 'heading' in section:
            doc.add_heading(section['heading'], level=section.get('level', 1))
        
        if 'text' in section:
            doc.add_paragraph(section['text'])
        
        if 'bullets' in section:
            for bullet in section['bullets']:
                doc.add_paragraph(bullet, style='List Bullet')
    
    def _generate_powerpoint(self, title: str, content: Any, **kwargs) -> Path:
        """Generate PowerPoint presentation."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError("python-pptx required. Install: pip install python-pptx")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        
        # Content slides
        if isinstance(content, list):
            for item in content:
                if isinstance(item, dict):
                    self._add_slide_to_pptx(prs, item)
                else:
                    self._add_text_slide(prs, "Content", str(item))
        elif isinstance(content, dict):
            for slide_title, slide_content in content.items():
                self._add_slide_to_pptx(prs, {
                    'title': slide_title,
                    'content': slide_content
                })
        
        # Save
        file_path = self.output_dir / f"{self._sanitize_filename(title)}.pptx"
        prs.save(str(file_path))
        
        logger.info(f"Generated PowerPoint: {file_path}")
        return file_path
    
    def _add_slide_to_pptx(self, prs, slide_data: Dict):
        """Add a slide to PowerPoint."""
        slide_title = slide_data.get('title', 'Slide')
        slide_content = slide_data.get('content', [])
        
        # Use bullet layout
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_title
        
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        if isinstance(slide_content, list):
            for item in slide_content:
                p = text_frame.add_paragraph()
                p.text = str(item)
                p.level = 0
        elif isinstance(slide_content, str):
            text_frame.text = slide_content
    
    def _add_text_slide(self, prs, title: str, text: str):
        """Add simple text slide."""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        
        content_top = Inches(2)
        content_height = Inches(4)
        content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
        content_frame = content_box.text_frame
        content_frame.text = text
    
    def _generate_pdf(self, title: str, content: Any, **kwargs) -> Path:
        """Generate PDF document."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
        except ImportError:
            raise ImportError("reportlab required. Install: pip install reportlab")
        
        file_path = self.output_dir / f"{self._sanitize_filename(title)}.pdf"
        
        doc = SimpleDocTemplate(str(file_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Title
        story.append(Paragraph(title, styles['Title']))
        story.append(Spacer(1, 12))
        
        # Content
        if isinstance(content, dict):
            for section_title, section_content in content.items():
                story.append(Paragraph(section_title, styles['Heading1']))
                if isinstance(section_content, list):
                    for item in section_content:
                        story.append(Paragraph(f"• {item}", styles['Normal']))
                else:
                    story.append(Paragraph(str(section_content), styles['Normal']))
                story.append(Spacer(1, 12))
        else:
            story.append(Paragraph(str(content), styles['Normal']))
        
        doc.build(story)
        
        logger.info(f"Generated PDF: {file_path}")
        return file_path
    
    def _generate_text(self, title: str, content: Any, **kwargs) -> Path:
        """Generate plain text file."""
        file_path = self.output_dir / f"{self._sanitize_filename(title)}.txt"
        
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(f"{title}\n")
            f.write("=" * len(title) + "\n\n")
            
            if isinstance(content, dict):
                for section, text in content.items():
                    f.write(f"\n{section}\n")
                    f.write("-" * len(section) + "\n")
                    if isinstance(text, list):
                        for item in text:
                            f.write(f"• {item}\n")
                    else:
                        f.write(f"{text}\n")
            else:
                f.write(str(content))
        
        logger.info(f"Generated text file: {file_path}")
        return file_path
    
    def _generate_markdown(self, title: str, content: Any, **kwargs) -> Path:
        """Generate Markdown file."""
        file_path = self.output_dir / f"{self._sanitize_filename(title)}.md"
        
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(f"# {title}\n\n")
            
            if isinstance(content, dict):
                for section, text in content.items():
                    f.write(f"## {section}\n\n")
                    if isinstance(text, list):
                        for item in text:
                            f.write(f"- {item}\n")
                        f.write("\n")
                    else:
                        f.write(f"{text}\n\n")
            else:
                f.write(str(content))
        
        logger.info(f"Generated markdown: {file_path}")
        return file_path
    
    def _sanitize_filename(self, filename: str) -> str:
        """Sanitize filename for filesystem."""
        import re
        # Remove invalid characters
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        # Limit length
        filename = filename[:200]
        return filename.strip()