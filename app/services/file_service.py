# app/services/file_service.py
"""
File service for handling uploads, downloads, and content extraction.
Files are stored in Redis (binary) with metadata.
Content is extracted for AI processing.
"""

import base64
import hashlib
import io
import json
import mimetypes
import os
import tempfile
import uuid
from dataclasses import dataclass, asdict
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple, BinaryIO

from app.utils.logger import setup_logger

logger = setup_logger(__name__)


@dataclass
class FileMetadata:
    """Metadata for an uploaded file."""
    file_id: str
    original_name: str
    mime_type: str
    size_bytes: int
    upload_time: str
    chat_id: Optional[str] = None
    user_id: Optional[int] = None
    checksum: str = ""
    extracted_content: str = ""
    content_type: str = "unknown"  # pdf, docx, xlsx, txt, image, etc.
    
    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> "FileMetadata":
        return cls(**{k: v for k, v in data.items() if k in cls.__dataclass_fields__})


class ContentExtractor:
    """Extracts text content from various file types."""
    
    SUPPORTED_TYPES = {
        'text': ['.txt', '.md', '.text', '.log', '.csv', '.json', '.xml', '.html', '.yaml', '.yml'],
        'pdf': ['.pdf'],
        'docx': ['.docx'],
        'xlsx': ['.xlsx', '.xls'],
        'pptx': ['.pptx', '.ppt'],
        'image': ['.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'],
    }
    
    @classmethod
    def get_content_type(cls, filename: str) -> str:
        """Determine content type from filename."""
        ext = Path(filename).suffix.lower()
        for content_type, extensions in cls.SUPPORTED_TYPES.items():
            if ext in extensions:
                return content_type
        return "unknown"
    
    @classmethod
    def extract(cls, file_data: bytes, filename: str) -> Tuple[str, str]:
        """
        Extract text content from file.
        Returns (content, content_type).
        """
        content_type = cls.get_content_type(filename)
        
        try:
            if content_type == 'text':
                return cls._extract_text(file_data), content_type
            elif content_type == 'pdf':
                return cls._extract_pdf(file_data), content_type
            elif content_type == 'docx':
                return cls._extract_docx(file_data), content_type
            elif content_type == 'xlsx':
                return cls._extract_xlsx(file_data), content_type
            elif content_type == 'pptx':
                return cls._extract_pptx(file_data), content_type
            elif content_type == 'image':
                return f"[Image: {filename}]", content_type
            else:
                return f"[Binary file: {filename}]", content_type
        except Exception as e:
            logger.error(f"Content extraction failed for {filename}: {e}")
            return f"[Extraction failed: {str(e)}]", content_type
    
    @classmethod
    def _extract_text(cls, data: bytes) -> str:
        """Extract from plain text files."""
        for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252']:
            try:
                return data.decode(encoding)
            except UnicodeDecodeError:
                continue
        return data.decode('utf-8', errors='replace')
    
    @classmethod
    def _extract_pdf(cls, data: bytes) -> str:
        """Extract text from PDF."""
        try:
            import pdfplumber
            with io.BytesIO(data) as pdf_file:
                with pdfplumber.open(pdf_file) as pdf:
                    text_parts = []
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
                    return '\n\n'.join(text_parts)
        except ImportError:
            # Fallback to pypdf
            try:
                from pypdf import PdfReader
                with io.BytesIO(data) as pdf_file:
                    reader = PdfReader(pdf_file)
                    return '\n\n'.join(
                        page.extract_text() or '' 
                        for page in reader.pages
                    )
            except ImportError:
                return "[PDF extraction requires pdfplumber or pypdf]"
    
    @classmethod
    def _extract_docx(cls, data: bytes) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document
            with io.BytesIO(data) as docx_file:
                doc = Document(docx_file)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return '\n\n'.join(paragraphs)
        except ImportError:
            return "[DOCX extraction requires python-docx]"
    
    @classmethod
    def _extract_xlsx(cls, data: bytes) -> str:
        """Extract text from Excel spreadsheet."""
        try:
            import pandas as pd
            with io.BytesIO(data) as xlsx_file:
                df = pd.read_excel(xlsx_file, sheet_name=None)
                parts = []
                for sheet_name, sheet_df in df.items():
                    parts.append(f"=== Sheet: {sheet_name} ===")
                    parts.append(sheet_df.to_string(index=False))
                return '\n\n'.join(parts)
        except ImportError:
            return "[XLSX extraction requires pandas and openpyxl]"
    
    @classmethod
    def _extract_pptx(cls, data: bytes) -> str:
        """Extract text from PowerPoint."""
        try:
            from pptx import Presentation
            with io.BytesIO(data) as pptx_file:
                prs = Presentation(pptx_file)
                text_parts = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_text = [f"=== Slide {i} ==="]
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                    text_parts.append('\n'.join(slide_text))
                return '\n\n'.join(text_parts)
        except ImportError:
            return "[PPTX extraction requires python-pptx]"


class FileService:
    """
    Service for file upload/download operations.
    Uses Redis for file storage (not SQLite).
    """
    
    # Redis key prefixes
    KEY_FILE_DATA = "file:data:"
    KEY_FILE_META = "file:meta:"
    KEY_CHAT_FILES = "chat:files:"
    
    # Limits
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    FILE_TTL = 86400 * 7  # 7 days
    
    def __init__(self):
        self._redis = None
        self._init_redis()
    
    def _init_redis(self):
        """Initialize Redis connection."""
        try:
            import redis
            self._redis = redis.Redis(
                host='localhost',
                port=6379,
                decode_responses=False,  # Need binary for file data
                socket_connect_timeout=2
            )
            self._redis.ping()
            logger.info("✅ FileService: Redis connected")
        except Exception as e:
            logger.warning(f"⚠️ FileService: Redis unavailable: {e}")
            self._redis = None
    
    @property
    def is_available(self) -> bool:
        """Check if file storage is available."""
        return self._redis is not None
    
    def upload_file(
        self,
        file_data: bytes,
        filename: str,
        chat_id: Optional[str] = None,
        user_id: Optional[int] = None
    ) -> Optional[FileMetadata]:
        """
        Upload a file to Redis storage.
        
        Returns FileMetadata with extracted content.
        """
        if not self.is_available:
            raise RuntimeError("File storage unavailable")
        
        if len(file_data) > self.MAX_FILE_SIZE:
            raise ValueError(f"File too large. Max size: {self.MAX_FILE_SIZE // (1024*1024)}MB")
        
        # Generate unique ID
        file_id = str(uuid.uuid4())
        
        # Calculate checksum
        checksum = hashlib.sha256(file_data).hexdigest()[:16]
        
        # Determine MIME type
        mime_type, _ = mimetypes.guess_type(filename)
        mime_type = mime_type or 'application/octet-stream'
        
        # Extract content
        extracted_content, content_type = ContentExtractor.extract(file_data, filename)
        
        # Create metadata
        metadata = FileMetadata(
            file_id=file_id,
            original_name=filename,
            mime_type=mime_type,
            size_bytes=len(file_data),
            upload_time=datetime.utcnow().isoformat(),
            chat_id=chat_id,
            user_id=user_id,
            checksum=checksum,
            extracted_content=extracted_content[:50000],  # Limit content size
            content_type=content_type
        )
        
        try:
            pipe = self._redis.pipeline()
            
            # Store file data (binary)
            pipe.setex(
                f"{self.KEY_FILE_DATA}{file_id}",
                self.FILE_TTL,
                file_data
            )
            
            # Store metadata (JSON)
            pipe.setex(
                f"{self.KEY_FILE_META}{file_id}",
                self.FILE_TTL,
                json.dumps(metadata.to_dict())
            )
            
            # Link to chat if provided
            if chat_id:
                pipe.sadd(f"{self.KEY_CHAT_FILES}{chat_id}", file_id)
                pipe.expire(f"{self.KEY_CHAT_FILES}{chat_id}", self.FILE_TTL)
            
            pipe.execute()
            
            logger.info(f"📁 File uploaded: {filename} ({file_id})")
            return metadata
            
        except Exception as e:
            logger.error(f"File upload failed: {e}")
            raise
    
    def get_file(self, file_id: str) -> Optional[Tuple[bytes, FileMetadata]]:
        """
        Retrieve a file by ID.
        
        Returns (file_data, metadata) or None if not found.
        """
        if not self.is_available:
            return None
        
        try:
            # Get metadata
            meta_json = self._redis.get(f"{self.KEY_FILE_META}{file_id}")
            if not meta_json:
                return None
            
            metadata = FileMetadata.from_dict(json.loads(meta_json))
            
            # Get file data
            file_data = self._redis.get(f"{self.KEY_FILE_DATA}{file_id}")
            if not file_data:
                return None
            
            return file_data, metadata
            
        except Exception as e:
            logger.error(f"File retrieval failed: {e}")
            return None
    
    def get_metadata(self, file_id: str) -> Optional[FileMetadata]:
        """Get file metadata only (no binary data)."""
        if not self.is_available:
            return None
        
        try:
            meta_json = self._redis.get(f"{self.KEY_FILE_META}{file_id}")
            if meta_json:
                return FileMetadata.from_dict(json.loads(meta_json))
            return None
        except Exception as e:
            logger.error(f"Metadata retrieval failed: {e}")
            return None
    
    def get_chat_files(self, chat_id: str) -> List[FileMetadata]:
        """Get all file metadata for a chat."""
        if not self.is_available:
            return []
        
        try:
            file_ids = self._redis.smembers(f"{self.KEY_CHAT_FILES}{chat_id}")
            files = []
            for file_id in file_ids:
                if isinstance(file_id, bytes):
                    file_id = file_id.decode()
                meta = self.get_metadata(file_id)
                if meta:
                    files.append(meta)
            return sorted(files, key=lambda f: f.upload_time, reverse=True)
        except Exception as e:
            logger.error(f"Chat files retrieval failed: {e}")
            return []
    
    def delete_file(self, file_id: str) -> bool:
        """Delete a file from storage."""
        if not self.is_available:
            return False
        
        try:
            pipe = self._redis.pipeline()
            pipe.delete(f"{self.KEY_FILE_DATA}{file_id}")
            pipe.delete(f"{self.KEY_FILE_META}{file_id}")
            results = pipe.execute()
            return any(results)
        except Exception as e:
            logger.error(f"File deletion failed: {e}")
            return False
    
    def get_extracted_content(self, file_id: str) -> Optional[str]:
        """Get just the extracted content for a file."""
        meta = self.get_metadata(file_id)
        return meta.extracted_content if meta else None


# Global instance
file_service = FileService()